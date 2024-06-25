from pptx import Presentation
from pptx.util import Inches

# Crear una presentación
prs = Presentation()

# Título de la diapositiva
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Sistema de Bombeo - Análisis y Ubicación de Bombas"
subtitle.text = "Detalles del perfil del terreno y la distribución de las bombas"

# Datos del perfil del terreno
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Perfil del Terreno (Distancia en metros, Altura en metros)"
content = slide.placeholders[1]
perfil_terreno = [
    (0, 375),
    (1100, 390),
    (2300, 420),
    (3100, 465),
    (4000, 507)
]
for punto in perfil_terreno:
    content.text += f"{punto}\n"

# Pérdidas de carga por segmento
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Pérdidas de Carga por Segmento"
content = slide.placeholders[1]
perdidas_carga_segmentos_actualizadas = [
    (1100, 15, 25.3, 40.3),
    (1200, 30, 27.6, 57.6),
    (800, 45, 18.4, 63.4),
    (900, 42, 20.7, 62.7)
]
for i, segmento in enumerate(perdidas_carga_segmentos_actualizadas):
    content.text += f"Segmento {i+1}: Longitud = {segmento[0]} m, Desnivel = {segmento[1]} m, Pérdida de fricción = {segmento[2]} m, Pérdida total = {segmento[3]} m\n"

# Ubicación de las bombas
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Ubicación de las Bombas"
content = slide.placeholders[1]
ubicaciones_bombas = [
    ("Primera Bomba", 0),
    ("Segunda Bomba", 2300),
    ("Tercera Bomba", 3100)
]
for bomba in ubicaciones_bombas:
    content.text += f"{bomba[0]}: {bomba[1]} metros\n"

# Guardar la presentación
pptx_output_path = "sistema_de_bombeo.pptx"
prs.save(pptx_output_path)
